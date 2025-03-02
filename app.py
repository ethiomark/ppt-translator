import streamlit as st
from pptx import Presentation
from deep_translator import GoogleTranslator
import io

st.title("PowerPoint Translator")

def translate_text(text, dest_language):
    """Translate text using Deep Translator."""
    return GoogleTranslator(source="auto", target=dest_language).translate(text)

def process_pptx_bytes(in_bytes):
    """Reads a PPTX from memory, translates notes, and returns updated file bytes."""
    prs = Presentation(io.BytesIO(in_bytes))

    # Your existing translation logic, simplified for brevity:
    for slide in prs.slides:
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                # Translate to Amharic & Swedish
                amharic = translate_text(notes_text, 'am')
                swedish = translate_text(notes_text, 'sv')

                # Replace the notes or add a shape, etc. 
                # (In this example, we just replace the notes text for simplicity.)
                # If you want to add textboxes, do that as in your script.
                slide.notes_slide.notes_text_frame.text = (
                    f"Amharic:\n{amharic}\n\nSwedish:\n{swedish}"
                )

    # Save the updated PPTX to memory
    out_bytes = io.BytesIO()
    prs.save(out_bytes)
    out_bytes.seek(0)
    return out_bytes

# Streamlit UI
uploaded_file = st.file_uploader("Upload a PowerPoint file (.pptx)", type=["pptx"])
if uploaded_file is not None:
    # Convert the uploaded file to bytes
    file_bytes = uploaded_file.read()

    # Process the PPTX in memory
    translated_pptx_bytes = process_pptx_bytes(file_bytes)

    # Download button
    st.download_button(
        label="Download Translated PPTX",
        data=translated_pptx_bytes,
        file_name="translated_output.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
